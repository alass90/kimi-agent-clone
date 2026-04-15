"""
Tool Executors — Complete Implementation
All 29 tools from Kimi's OK Computer agent, faithfully reproduced.

Architecture:
  - E2B Sandbox: Cloud-based code execution, shell, and filesystem
  - Playwright Browser: Full browser automation with anti-detection
  - OpenAI: Image generation (DALL-E), Speech (TTS)
  - Data Sources: Yahoo Finance, World Bank, arXiv, Google Scholar
  - Asset Extraction: Image bbox detection and cropping via Pillow/OpenCV
  - Slides: HTML → PPTX conversion
  - Sound Effects: Text-to-sound via API
  - Todo: Persistent task management (.todo.jsonl)
"""
import asyncio
import json
import os
import subprocess
import shutil
import base64
import logging
import signal
import time
import re
import urllib.parse
from pathlib import Path
from typing import Optional, Dict, Any, List, Tuple
from datetime import datetime

logger = logging.getLogger(__name__)

# ─── Workspace paths ─────────────────────────────────────────────────────
WORKSPACE = Path(os.getenv("WORKSPACE_DIR", "/home/ubuntu/kimi-clone/workspace"))
OUTPUT_DIR = WORKSPACE / "output"
UPLOAD_DIR = WORKSPACE / "upload"
DEPLOY_DIR = WORKSPACE / "deploy"
TODO_FILE = WORKSPACE / ".todo.jsonl"

for d in [OUTPUT_DIR, UPLOAD_DIR, DEPLOY_DIR]:
    d.mkdir(parents=True, exist_ok=True)

# ─── Session State ────────────────────────────────────────────────────────
_files_read: set = set()
_e2b_sandbox = None
_browser_instance = None
_browser_page = None


# ═══════════════════════════════════════════════════════════════════════════
# E2B Sandbox Manager (Cloud Code Execution)
# ═══════════════════════════════════════════════════════════════════════════

class E2BSandboxManager:
    """
    Manages an E2B cloud sandbox for secure code execution and shell commands.
    Falls back to local execution if E2B is not configured.
    """

    def __init__(self):
        self.sandbox = None
        self.use_e2b = bool(os.getenv("E2B_API_KEY"))

    async def get_sandbox(self):
        """Get or create an E2B sandbox instance."""
        if not self.use_e2b:
            return None
        if self.sandbox is None:
            try:
                from e2b_code_interpreter import AsyncSandbox
                self.sandbox = await AsyncSandbox.create(timeout=300)
                logger.info("E2B sandbox created successfully")
                # Initialize with common packages
                await self.sandbox.run_code(
                    "import matplotlib\nmatplotlib.use('Agg')\n"
                    "import matplotlib.pyplot as plt\nimport numpy as np\n"
                    "import pandas as pd\nfrom IPython.display import display\n"
                    "plt.style.use('default')\n"
                    "plt.rcParams['figure.figsize'] = [8, 6]\n"
                    "plt.rcParams['figure.dpi'] = 100\n"
                    "plt.rcParams['axes.grid'] = True\n"
                    "plt.rcParams['figure.facecolor'] = 'white'\n"
                )
            except Exception as e:
                logger.warning(f"E2B sandbox creation failed: {e}. Falling back to local.")
                self.use_e2b = False
                self.sandbox = None
        return self.sandbox

    async def run_code(self, code: str, timeout: int = 30) -> Dict[str, Any]:
        """Execute Python code in E2B sandbox or locally."""
        sandbox = await self.get_sandbox()
        if sandbox:
            return await self._run_e2b(sandbox, code, timeout)
        return self._run_local(code, timeout)

    async def _run_e2b(self, sandbox, code: str, timeout: int) -> Dict[str, Any]:
        """Execute code in E2B cloud sandbox."""
        try:
            execution = await sandbox.run_code(code, timeout=timeout)
            output_parts = []
            images = []
            error = None

            # Collect stdout
            if execution.logs and execution.logs.stdout:
                for msg in execution.logs.stdout:
                    output_parts.append(msg.line if hasattr(msg, 'line') else str(msg))

            # Collect stderr
            if execution.logs and execution.logs.stderr:
                for msg in execution.logs.stderr:
                    output_parts.append(f"[stderr] {msg.line if hasattr(msg, 'line') else str(msg)}")

            # Collect results (text, images, etc.)
            if execution.results:
                for result in execution.results:
                    text_repr = str(result)
                    if text_repr:
                        output_parts.append(text_repr)
                    # Check for image data
                    png = result._repr_png_() if hasattr(result, '_repr_png_') else None
                    if png:
                        images.append(png)

            # Check for errors
            if execution.error:
                error = f"{execution.error.name}: {execution.error.value}\n{execution.error.traceback}"

            final_output = "\n".join(output_parts).strip()
            if len(final_output) > 10000:
                final_output = final_output[:10000] + "\n... [output truncated]"

            return {
                "success": error is None,
                "output": final_output,
                "error": error,
                "images": images,
                "runtime": "e2b",
            }
        except Exception as e:
            return {
                "success": False,
                "output": "",
                "error": f"E2B execution error: {str(e)}",
                "images": [],
                "runtime": "e2b",
            }

    def _run_local(self, code: str, timeout: int) -> Dict[str, Any]:
        """Fallback: execute code locally via Jupyter kernel."""
        kernel = _get_kernel()
        result = kernel.execute(code, timeout=timeout)
        result["runtime"] = "local"
        return result

    async def run_shell(self, command: str, timeout: int = 600) -> Dict[str, Any]:
        """Execute shell command in E2B sandbox or locally."""
        sandbox = await self.get_sandbox()
        if sandbox:
            try:
                result = sandbox.commands.run(command, timeout=timeout)
                output = (result.stdout or "") + (result.stderr or "")
                if len(output) > 10000:
                    output = output[:10000] + "\n... [output truncated]"
                return {
                    "success": result.exit_code == 0,
                    "output": output,
                    "return_code": result.exit_code,
                    "runtime": "e2b",
                }
            except Exception as e:
                return {"error": f"E2B shell error: {str(e)}", "runtime": "e2b"}

        # Local fallback
        try:
            result = subprocess.run(
                command, shell=True, capture_output=True, text=True,
                timeout=timeout, cwd=str(WORKSPACE),
            )
            output = result.stdout + result.stderr
            if len(output) > 10000:
                output = output[:10000] + "\n... [output truncated]"
            return {
                "success": result.returncode == 0,
                "output": output,
                "return_code": result.returncode,
                "runtime": "local",
            }
        except subprocess.TimeoutExpired:
            return {"error": f"Command timed out after {timeout}s"}
        except Exception as e:
            return {"error": str(e)}

    async def read_file(self, path: str) -> Optional[str]:
        """Read file from E2B sandbox or local filesystem."""
        sandbox = await self.get_sandbox()
        if sandbox:
            try:
                content = await sandbox.files.read(path)
                return content
            except Exception:
                pass
        # Local fallback
        try:
            with open(path, "r", encoding="utf-8", errors="replace") as f:
                return f.read()
        except Exception:
            return None

    async def write_file(self, path: str, content: str) -> bool:
        """Write file to E2B sandbox or local filesystem."""
        sandbox = await self.get_sandbox()
        if sandbox:
            try:
                await sandbox.files.write(path, content)
                return True
            except Exception:
                pass
        # Local fallback
        try:
            os.makedirs(os.path.dirname(path), exist_ok=True)
            with open(path, "w", encoding="utf-8") as f:
                f.write(content)
            return True
        except Exception:
            return False

    async def list_files(self, path: str) -> List[str]:
        """List files in E2B sandbox or local filesystem."""
        sandbox = await self.get_sandbox()
        if sandbox:
            try:
                files = await sandbox.files.list(path)
                return [f.name for f in files]
            except Exception:
                pass
        try:
            return os.listdir(path)
        except Exception:
            return []

    async def shutdown(self):
        """Shutdown the E2B sandbox."""
        if self.sandbox:
            try:
                await self.sandbox.kill()
            except Exception:
                pass
            self.sandbox = None


# Global sandbox manager
_sandbox_mgr = E2BSandboxManager()


# ═══════════════════════════════════════════════════════════════════════════
# IPython / Jupyter Kernel (Local Fallback — mirrors jupyter_kernel.py)
# ═══════════════════════════════════════════════════════════════════════════

class IPythonKernel:
    """
    Robust wrapper around jupyter_client for local code execution.
    Mirrors Kimi's jupyter_kernel.py with:
    - PID tracking for process management
    - SIGINT support for interrupting long-running code
    - Automatic recovery on kernel death
    - Output truncation at 10000 chars
    """

    def __init__(self):
        self.km = None
        self.kc = None
        self.kernel_pid = None
        self._start_kernel()

    def _start_kernel(self):
        try:
            from jupyter_client.manager import KernelManager
            if self.km:
                self.shutdown()
            self.km = KernelManager()
            self.km.start_kernel()
            self.kc = self.km.client()
            self.kc.start_channels()
            self.kc.wait_for_ready(timeout=30)
            # Track kernel PID for SIGINT
            try:
                self.kernel_pid = self.km.kernel.pid
            except Exception:
                self.kernel_pid = None
            # Initialize environment
            init_code = (
                "import matplotlib\nmatplotlib.use('Agg')\n"
                "import matplotlib.pyplot as plt\nimport numpy as np\n"
                "import pandas as pd\nfrom IPython.display import display\n"
                "%matplotlib inline\n"
                "plt.style.use('default')\n"
                "plt.rcParams['figure.figsize'] = [8, 6]\n"
                "plt.rcParams['figure.dpi'] = 100\n"
                "plt.rcParams['axes.grid'] = True\n"
                "plt.rcParams['figure.facecolor'] = 'white'\n"
            )
            self.execute(init_code, timeout=15)
            logger.info(f"IPython kernel initialized (PID: {self.kernel_pid})")
        except Exception as e:
            logger.error(f"Kernel init error: {e}")
            raise

    def interrupt(self):
        """Send SIGINT to the kernel process (mirrors Kimi's interrupt capability)."""
        if self.kernel_pid:
            try:
                os.kill(self.kernel_pid, signal.SIGINT)
                logger.info(f"SIGINT sent to kernel PID {self.kernel_pid}")
            except ProcessLookupError:
                logger.warning("Kernel process not found for interrupt")
        elif self.km:
            try:
                self.km.interrupt_kernel()
            except Exception:
                pass

    def execute(self, code: str, timeout: int = 30) -> Dict[str, Any]:
        try:
            if not self.kc or not self.km:
                self._start_kernel()
            if not self.km.is_alive():
                self._start_kernel()

            msg_id = self.kc.execute(code)
            output_parts = []
            error = None
            images = []
            start = time.time()

            while True:
                if time.time() - start > timeout:
                    self.interrupt()
                    return {
                        "success": False,
                        "output": "".join(output_parts).strip(),
                        "error": f"Execution timed out after {timeout}s (interrupted)",
                        "images": images,
                    }
                try:
                    msg = self.kc.get_iopub_msg(timeout=2)
                    msg_type = msg["header"]["msg_type"]

                    if msg_type == "stream":
                        output_parts.append(msg["content"]["text"])
                    elif msg_type == "error":
                        error = "\n".join(msg["content"]["traceback"])
                    elif msg_type == "execute_result":
                        data = msg["content"].get("data", {})
                        if "text/plain" in data:
                            output_parts.append(data["text/plain"])
                        if "image/png" in data:
                            images.append(data["image/png"])
                    elif msg_type == "display_data":
                        data = msg["content"].get("data", {})
                        if "image/png" in data:
                            images.append(data["image/png"])
                        if "text/html" in data:
                            output_parts.append(data["text/html"])
                        elif "text/plain" in data:
                            output_parts.append(data["text/plain"])

                    if (
                        msg["parent_header"].get("msg_id") == msg_id
                        and msg_type == "status"
                        and msg["content"]["execution_state"] == "idle"
                    ):
                        break
                except Exception as e:
                    if "Timeout" in str(type(e).__name__) or "Empty" in str(type(e).__name__):
                        continue
                    raise

            final_output = "".join(output_parts).strip()
            if len(final_output) > 10000:
                final_output = final_output[:10000] + "\n... [output truncated]"

            return {
                "success": error is None,
                "output": final_output,
                "error": error,
                "images": images,
            }
        except Exception as e:
            logger.error(f"Execution error: {e}")
            try:
                self._start_kernel()
            except Exception:
                pass
            return {
                "success": False,
                "output": "",
                "error": f"{e.__class__.__name__}: {str(e)}",
                "images": [],
            }

    def restart(self):
        self.shutdown()
        self._start_kernel()

    def shutdown(self):
        try:
            if self.kc:
                self.kc.stop_channels()
            if self.km:
                self.km.shutdown_kernel(now=True)
        except Exception:
            pass
        self.km = None
        self.kc = None
        self.kernel_pid = None


_ipython_kernel = None

def _get_kernel() -> IPythonKernel:
    global _ipython_kernel
    if _ipython_kernel is None:
        _ipython_kernel = IPythonKernel()
    return _ipython_kernel


# ═══════════════════════════════════════════════════════════════════════════
# Browser Automation (Playwright — mirrors browser_guard.py)
# ═══════════════════════════════════════════════════════════════════════════

class PlaywrightBrowser:
    """
    Full browser automation with Playwright.
    Mirrors Kimi's browser_guard.py with:
    - Chrome in stealth mode (anti-detection)
    - Persistent context (cookies, localStorage)
    - Element indexing for LLM interaction
    - Screenshot capture
    - Scroll management
    - Text search
    """

    def __init__(self):
        self.playwright = None
        self.browser = None
        self.context = None
        self.page = None
        self.elements = []
        self.current_url = None

    async def _ensure_browser(self):
        """Launch browser if not already running."""
        if self.page and not self.page.is_closed():
            return

        try:
            from playwright.async_api import async_playwright

            if not self.playwright:
                self.playwright = await async_playwright().start()

            # Launch with stealth settings (mirrors browser_guard.py)
            self.browser = await self.playwright.chromium.launch(
                headless=True,
                args=[
                    "--no-sandbox",
                    "--disable-setuid-sandbox",
                    "--disable-blink-features=AutomationControlled",
                    "--disable-infobars",
                    "--window-size=1920,1080",
                    "--disable-extensions",
                    "--disable-gpu",
                    "--disable-dev-shm-usage",
                ]
            )

            # Create context with stealth settings
            self.context = await self.browser.new_context(
                viewport={"width": 1920, "height": 1080},
                user_agent=(
                    "Mozilla/5.0 (Windows NT 10.0; Win64; x64) "
                    "AppleWebKit/537.36 (KHTML, like Gecko) "
                    "Chrome/120.0.0.0 Safari/537.36"
                ),
                java_script_enabled=True,
                bypass_csp=True,
                ignore_https_errors=True,
            )

            self.page = await self.context.new_page()

            # Inject stealth scripts (mirrors Kimi's anti-detection)
            await self.page.add_init_script("""
                // Remove webdriver flag
                Object.defineProperty(navigator, 'webdriver', { get: () => undefined });
                // Mock plugins
                Object.defineProperty(navigator, 'plugins', {
                    get: () => [1, 2, 3, 4, 5]
                });
                // Mock languages
                Object.defineProperty(navigator, 'languages', {
                    get: () => ['en-US', 'en']
                });
                // Override permissions
                const originalQuery = window.navigator.permissions.query;
                window.navigator.permissions.query = (parameters) =>
                    parameters.name === 'notifications'
                        ? Promise.resolve({ state: Notification.permission })
                        : originalQuery(parameters);
                // Chrome runtime mock
                window.chrome = { runtime: {} };
            """)

            logger.info("Playwright browser initialized with stealth mode")
        except ImportError:
            logger.warning("Playwright not installed. Browser tools will use curl fallback.")
            raise
        except Exception as e:
            logger.error(f"Browser init error: {e}")
            raise

    async def _extract_elements(self) -> List[Dict]:
        """Extract interactive elements from the page (mirrors Kimi's element indexing)."""
        if not self.page:
            return []

        try:
            elements = await self.page.evaluate("""
                () => {
                    const interactiveSelectors = 'a, button, input, select, textarea, [role="button"], [role="link"], [role="tab"], [onclick], [tabindex]';
                    const elements = document.querySelectorAll(interactiveSelectors);
                    return Array.from(elements).slice(0, 200).map((el, idx) => ({
                        index: idx,
                        tag: el.tagName.toLowerCase(),
                        text: (el.textContent || '').trim().substring(0, 100),
                        type: el.type || '',
                        href: el.href || '',
                        placeholder: el.placeholder || '',
                        value: el.value || '',
                        aria_label: el.getAttribute('aria-label') || '',
                        visible: el.offsetParent !== null,
                    }));
                }
            """)
            self.elements = elements
            return elements
        except Exception:
            return []

    async def visit(self, url: str, need_screenshot: bool = False) -> Dict:
        """Navigate to URL and extract page info."""
        try:
            await self._ensure_browser()
            response = await self.page.goto(url, wait_until="domcontentloaded", timeout=30000)
            await self.page.wait_for_timeout(2000)  # Wait for dynamic content
            self.current_url = self.page.url

            # Extract page content
            title = await self.page.title()
            content = await self.page.evaluate("() => document.body.innerText.substring(0, 10000)")
            elements = await self._extract_elements()

            result = {
                "success": True,
                "url": self.current_url,
                "title": title,
                "status": response.status if response else None,
                "content_preview": content[:3000],
                "interactive_elements": len(elements),
                "elements": [
                    f"{e['index']}[:{e['tag']}]{e['text'][:50]}"
                    for e in elements[:30] if e.get("visible", True)
                ],
            }

            if need_screenshot:
                screenshot_path = str(OUTPUT_DIR / f"screenshot_{int(time.time())}.png")
                await self.page.screenshot(path=screenshot_path, full_page=False)
                result["screenshot_path"] = screenshot_path

            return result
        except ImportError:
            # Fallback to curl
            return await self._curl_fallback(url)
        except Exception as e:
            return {"error": f"Browser visit failed: {str(e)}"}

    async def _curl_fallback(self, url: str) -> Dict:
        """Fallback when Playwright is not available."""
        try:
            result = subprocess.run(
                f'curl -sL -A "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36" "{url}" 2>/dev/null | head -c 20000',
                shell=True, capture_output=True, text=True, timeout=30,
            )
            self.current_url = url
            return {
                "success": True,
                "url": url,
                "content_preview": result.stdout[:3000],
                "content_length": len(result.stdout),
                "note": "Using curl fallback (Playwright not available)",
            }
        except Exception as e:
            return {"error": str(e)}

    async def click(self, element_index: int) -> Dict:
        """Click an element by index."""
        if not self.page:
            return {"error": "No page loaded. Use browser_visit first."}
        try:
            result = await self.page.evaluate(f"""
                () => {{
                    const elements = document.querySelectorAll('a, button, input, select, textarea, [role="button"], [role="link"], [role="tab"], [onclick], [tabindex]');
                    const el = elements[{element_index}];
                    if (!el) return {{ error: 'Element not found at index {element_index}' }};
                    el.click();
                    return {{ success: true, tag: el.tagName, text: (el.textContent || '').trim().substring(0, 100) }};
                }}
            """)
            await self.page.wait_for_timeout(1000)
            self.current_url = self.page.url
            return result
        except Exception as e:
            return {"error": str(e)}

    async def input_text(self, element_index: int, content: str) -> Dict:
        """Type text into a form field."""
        if not self.page:
            return {"error": "No page loaded. Use browser_visit first."}
        try:
            result = await self.page.evaluate(f"""
                () => {{
                    const elements = document.querySelectorAll('a, button, input, select, textarea, [role="button"], [role="link"], [role="tab"], [onclick], [tabindex]');
                    const el = elements[{element_index}];
                    if (!el) return {{ error: 'Element not found' }};
                    el.focus();
                    el.value = '';
                    return {{ success: true, tag: el.tagName }};
                }}
            """)
            if result.get("error"):
                return result
            # Use Playwright's type for realistic input
            selectors = 'a, button, input, select, textarea, [role="button"], [role="link"], [role="tab"], [onclick], [tabindex]'
            elements = await self.page.query_selector_all(selectors)
            if element_index < len(elements):
                await elements[element_index].fill(content)
            return {"success": True, "content": content}
        except Exception as e:
            return {"error": str(e)}

    async def scroll(self, direction: str = "down", amount: int = 500) -> Dict:
        """Scroll the page."""
        if not self.page:
            return {"error": "No page loaded."}
        try:
            delta = amount if direction == "down" else -amount
            await self.page.evaluate(f"window.scrollBy(0, {delta})")
            await self.page.wait_for_timeout(500)
            scroll_pos = await self.page.evaluate("() => window.scrollY")
            page_height = await self.page.evaluate("() => document.body.scrollHeight")
            return {
                "success": True,
                "direction": direction,
                "scroll_position": scroll_pos,
                "page_height": page_height,
            }
        except Exception as e:
            return {"error": str(e)}

    async def screenshot(self, download_path: str = None) -> Dict:
        """Take a screenshot."""
        if not self.page:
            return {"error": "No page loaded."}
        try:
            path = download_path or str(OUTPUT_DIR / f"screenshot_{int(time.time())}.png")
            await self.page.screenshot(path=path, full_page=False)
            return {"success": True, "path": path, "url": self.current_url}
        except Exception as e:
            return {"error": str(e)}

    async def find_text(self, keyword: str, skip: int = 0) -> Dict:
        """Search for text on the page."""
        if not self.page:
            return {"error": "No page loaded."}
        try:
            content = await self.page.evaluate("() => document.body.innerText")
            occurrences = []
            lower_content = content.lower()
            lower_keyword = keyword.lower()
            start = 0
            while True:
                idx = lower_content.find(lower_keyword, start)
                if idx == -1:
                    break
                context_start = max(0, idx - 50)
                context_end = min(len(content), idx + len(keyword) + 50)
                occurrences.append({
                    "position": idx,
                    "context": content[context_start:context_end],
                })
                start = idx + 1

            return {
                "found": len(occurrences) > 0,
                "keyword": keyword,
                "total_matches": len(occurrences),
                "matches": occurrences[skip:skip + 5],
            }
        except Exception as e:
            return {"error": str(e)}

    async def close(self):
        """Close the browser."""
        try:
            if self.page:
                await self.page.close()
            if self.context:
                await self.context.close()
            if self.browser:
                await self.browser.close()
            if self.playwright:
                await self.playwright.stop()
        except Exception:
            pass
        self.page = None
        self.context = None
        self.browser = None
        self.playwright = None


# Global browser instance
_browser = PlaywrightBrowser()


# ═══════════════════════════════════════════════════════════════════════════
# Data Source Clients (mirrors Kimi's get_data_source tool)
# ═══════════════════════════════════════════════════════════════════════════

class DataSourceClient:
    """
    Unified client for external data sources.
    Mirrors Kimi's datasource registry:
    - Yahoo Finance (yfinance)
    - World Bank Open Data (wbgapi)
    - arXiv (arxiv API)
    - Google Scholar (scholarly)
    """

    @staticmethod
    async def yahoo_finance(query: str, params: Dict = None) -> Dict:
        """Fetch financial data from Yahoo Finance."""
        try:
            code = f"""
import yfinance as yf
import json

ticker = yf.Ticker("{query}")
info = ticker.info
hist = ticker.history(period="1mo")

result = {{
    "symbol": "{query}",
    "name": info.get("longName", "N/A"),
    "price": info.get("currentPrice", info.get("regularMarketPrice", "N/A")),
    "currency": info.get("currency", "N/A"),
    "market_cap": info.get("marketCap", "N/A"),
    "pe_ratio": info.get("trailingPE", "N/A"),
    "52w_high": info.get("fiftyTwoWeekHigh", "N/A"),
    "52w_low": info.get("fiftyTwoWeekLow", "N/A"),
    "volume": info.get("volume", "N/A"),
    "sector": info.get("sector", "N/A"),
    "industry": info.get("industry", "N/A"),
    "history_1mo": hist.to_dict() if not hist.empty else {{}},
}}
print(json.dumps(result, default=str))
"""
            result = await _sandbox_mgr.run_code(code, timeout=30)
            if result.get("success") and result.get("output"):
                try:
                    return {"success": True, "data": json.loads(result["output"])}
                except json.JSONDecodeError:
                    return {"success": True, "data": result["output"]}
            return {"success": False, "error": result.get("error", "Unknown error")}
        except Exception as e:
            return {"error": str(e)}

    @staticmethod
    async def world_bank(indicator: str, country: str = "all", params: Dict = None) -> Dict:
        """Fetch data from World Bank Open Data."""
        try:
            url = f"https://api.worldbank.org/v2/country/{country}/indicator/{indicator}?format=json&per_page=50"
            result = subprocess.run(
                f'curl -s "{url}"', shell=True, capture_output=True, text=True, timeout=15,
            )
            if result.returncode == 0:
                data = json.loads(result.stdout)
                if isinstance(data, list) and len(data) > 1:
                    records = []
                    for item in data[1][:20]:
                        records.append({
                            "country": item.get("country", {}).get("value"),
                            "year": item.get("date"),
                            "value": item.get("value"),
                            "indicator": item.get("indicator", {}).get("value"),
                        })
                    return {"success": True, "data": records, "total": data[0].get("total", 0)}
            return {"success": False, "error": "No data found"}
        except Exception as e:
            return {"error": str(e)}

    @staticmethod
    async def arxiv_search(query: str, max_results: int = 10) -> Dict:
        """Search arXiv for academic papers."""
        try:
            encoded = urllib.parse.quote(query)
            url = f"http://export.arxiv.org/api/query?search_query=all:{encoded}&start=0&max_results={max_results}"
            result = subprocess.run(
                f'curl -s "{url}"', shell=True, capture_output=True, text=True, timeout=15,
            )
            if result.returncode == 0:
                # Parse XML response
                import xml.etree.ElementTree as ET
                root = ET.fromstring(result.stdout)
                ns = {"atom": "http://www.w3.org/2005/Atom"}
                papers = []
                for entry in root.findall("atom:entry", ns):
                    title = entry.find("atom:title", ns)
                    summary = entry.find("atom:summary", ns)
                    published = entry.find("atom:published", ns)
                    link = entry.find("atom:id", ns)
                    authors = entry.findall("atom:author/atom:name", ns)
                    papers.append({
                        "title": title.text.strip() if title is not None else "",
                        "summary": (summary.text.strip()[:300] + "...") if summary is not None else "",
                        "published": published.text if published is not None else "",
                        "url": link.text if link is not None else "",
                        "authors": [a.text for a in authors[:5]],
                    })
                return {"success": True, "papers": papers, "query": query}
            return {"success": False, "error": "arXiv API request failed"}
        except Exception as e:
            return {"error": str(e)}

    @staticmethod
    async def google_scholar(query: str, max_results: int = 10) -> Dict:
        """Search Google Scholar (via scraping fallback)."""
        try:
            encoded = urllib.parse.quote(query)
            url = f"https://scholar.google.com/scholar?q={encoded}&hl=en"
            result = subprocess.run(
                f'curl -s -A "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36" "{url}" 2>/dev/null | head -c 50000',
                shell=True, capture_output=True, text=True, timeout=15,
            )
            if result.returncode == 0:
                # Basic parsing of Google Scholar HTML
                from html.parser import HTMLParser
                titles = re.findall(r'<h3[^>]*class="gs_rt"[^>]*>(.*?)</h3>', result.stdout, re.DOTALL)
                snippets = re.findall(r'<div[^>]*class="gs_rs"[^>]*>(.*?)</div>', result.stdout, re.DOTALL)
                papers = []
                for i, title in enumerate(titles[:max_results]):
                    clean_title = re.sub(r'<[^>]+>', '', title).strip()
                    clean_snippet = re.sub(r'<[^>]+>', '', snippets[i]).strip() if i < len(snippets) else ""
                    papers.append({
                        "title": clean_title,
                        "snippet": clean_snippet[:200],
                    })
                return {"success": True, "papers": papers, "query": query}
            return {"success": False, "error": "Google Scholar request failed"}
        except Exception as e:
            return {"error": str(e)}


_data_client = DataSourceClient()


# ═══════════════════════════════════════════════════════════════════════════
# Web Search (DuckDuckGo + Google fallback)
# ═══════════════════════════════════════════════════════════════════════════

async def _perform_web_search(query: str, count: int = 5) -> Dict:
    """Perform web search using DuckDuckGo HTML API with parsing."""
    try:
        encoded = urllib.parse.quote(query)
        result = subprocess.run(
            f'curl -s -A "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36" '
            f'"https://html.duckduckgo.com/html/?q={encoded}" 2>/dev/null',
            shell=True, capture_output=True, text=True, timeout=15,
        )
        if result.returncode == 0 and result.stdout:
            # Parse DuckDuckGo HTML results
            results = []
            # Extract result blocks
            links = re.findall(r'<a[^>]*class="result__a"[^>]*href="([^"]*)"[^>]*>(.*?)</a>', result.stdout, re.DOTALL)
            snippets = re.findall(r'<a[^>]*class="result__snippet"[^>]*>(.*?)</a>', result.stdout, re.DOTALL)

            for i, (url, title) in enumerate(links[:count]):
                clean_title = re.sub(r'<[^>]+>', '', title).strip()
                clean_snippet = re.sub(r'<[^>]+>', '', snippets[i]).strip() if i < len(snippets) else ""
                # Decode DuckDuckGo redirect URL
                if "uddg=" in url:
                    try:
                        actual_url = urllib.parse.unquote(url.split("uddg=")[1].split("&")[0])
                    except Exception:
                        actual_url = url
                else:
                    actual_url = url
                results.append({
                    "title": clean_title,
                    "url": actual_url,
                    "snippet": clean_snippet[:200],
                })

            if results:
                return {"success": True, "query": query, "results": results}

        # Fallback: return raw content
        return {
            "success": True,
            "query": query,
            "results": result.stdout[:5000] if result.stdout else "No results found",
            "note": "Raw search results returned",
        }
    except Exception as e:
        return {"error": f"Search failed: {str(e)}"}


# ═══════════════════════════════════════════════════════════════════════════
# Image Generation (OpenAI DALL-E with full ratio/resolution support)
# ═══════════════════════════════════════════════════════════════════════════

RATIO_TO_SIZE = {
    "1:1": "1024x1024",
    "3:2": "1536x1024",
    "2:3": "1024x1536",
    "4:3": "1536x1024",
    "3:4": "1024x1536",
    "16:9": "1792x1024",
    "9:16": "1024x1792",
    "21:9": "1792x1024",
}


async def _generate_image_impl(prompt: str, output_path: str, ratio: str = "1:1", resolution: str = "1K") -> Dict:
    """Generate image using OpenAI DALL-E 3."""
    try:
        from openai import OpenAI
        client = OpenAI()
        size = RATIO_TO_SIZE.get(ratio, "1024x1024")
        quality = "hd" if resolution in ("2K", "4K") else "standard"

        response = client.images.generate(
            model="dall-e-3",
            prompt=prompt,
            size=size,
            quality=quality,
            n=1,
        )
        image_url = response.data[0].url
        revised_prompt = response.data[0].revised_prompt

        # Download the image
        import requests
        img_data = requests.get(image_url, timeout=30).content
        os.makedirs(os.path.dirname(output_path) if os.path.dirname(output_path) else ".", exist_ok=True)
        with open(output_path, "wb") as f:
            f.write(img_data)

        return {
            "success": True,
            "output_path": output_path,
            "prompt": prompt,
            "revised_prompt": revised_prompt,
            "size": size,
            "quality": quality,
            "bytes": len(img_data),
        }
    except Exception as e:
        return {"error": f"Image generation failed: {str(e)}"}


# ═══════════════════════════════════════════════════════════════════════════
# Speech / TTS (OpenAI TTS)
# ═══════════════════════════════════════════════════════════════════════════

async def _generate_speech_impl(text: str, output_path: str, voice: str = "alloy", speed: float = 1.0) -> Dict:
    """Generate speech from text using OpenAI TTS."""
    try:
        from openai import OpenAI
        client = OpenAI()
        response = client.audio.speech.create(
            model="tts-1",
            voice=voice,
            input=text,
            speed=speed,
        )
        os.makedirs(os.path.dirname(output_path) if os.path.dirname(output_path) else ".", exist_ok=True)
        response.stream_to_file(output_path)
        file_size = os.path.getsize(output_path)
        return {
            "success": True,
            "output_path": output_path,
            "voice": voice,
            "speed": speed,
            "text_length": len(text),
            "file_size_bytes": file_size,
        }
    except Exception as e:
        return {"error": f"Speech generation failed: {str(e)}"}


# ═══════════════════════════════════════════════════════════════════════════
# Sound Effects (via code generation — synthesize with pydub/numpy)
# ═══════════════════════════════════════════════════════════════════════════

async def _generate_sound_effect_impl(prompt: str, output_path: str, duration: float = 3.0) -> Dict:
    """Generate sound effects using synthesis (numpy + scipy)."""
    try:
        # Escape prompt for inclusion in code string
        safe_prompt = prompt.replace('"', '\\"')
        code = f"""
import numpy as np
from scipy.io import wavfile
import os

# Generate sound effect based on prompt
sr = 44100
duration = {duration}
t = np.linspace(0, duration, int(sr * duration), endpoint=False)

# Create a layered sound effect
prompt = "{safe_prompt}"
if any(w in prompt.lower() for w in ['bell', 'ding', 'chime', 'notification']):
    freq = 880
    signal = np.sin(2 * np.pi * freq * t) * np.exp(-3 * t)
    signal += 0.3 * np.sin(2 * np.pi * freq * 2 * t) * np.exp(-4 * t)
elif any(w in prompt.lower() for w in ['whoosh', 'wind', 'sweep']):
    noise = np.random.randn(len(t))
    freq_sweep = np.linspace(200, 2000, len(t))
    signal = noise * np.sin(2 * np.pi * freq_sweep * t / sr) * np.exp(-t)
elif any(w in prompt.lower() for w in ['click', 'tap', 'button']):
    signal = np.zeros(len(t))
    click_len = int(0.01 * sr)
    signal[:click_len] = np.sin(2 * np.pi * 1000 * t[:click_len]) * np.exp(-50 * t[:click_len])
elif any(w in prompt.lower() for w in ['alarm', 'siren', 'warning']):
    freq = 440 + 220 * np.sin(2 * np.pi * 3 * t)
    signal = np.sin(2 * np.pi * freq * t)
elif any(w in prompt.lower() for w in ['rain', 'water', 'ocean']):
    signal = np.random.randn(len(t)) * 0.3
    # Low-pass filter approximation
    for i in range(1, len(signal)):
        signal[i] = 0.99 * signal[i-1] + 0.01 * signal[i]
else:
    # Generic tone
    signal = np.sin(2 * np.pi * 440 * t) * np.exp(-t / duration)
    signal += 0.5 * np.sin(2 * np.pi * 660 * t) * np.exp(-t / (duration * 0.7))

# Normalize
signal = signal / np.max(np.abs(signal)) * 0.8
signal_int16 = (signal * 32767).astype(np.int16)

output_path = "{output_path}"
os.makedirs(os.path.dirname(output_path) if os.path.dirname(output_path) else ".", exist_ok=True)
wavfile.write(output_path, sr, signal_int16)
print(f"Sound effect saved: {{os.path.getsize(output_path)}} bytes")
"""
        result = await _sandbox_mgr.run_code(code, timeout=15)
        if result.get("success"):
            return {
                "success": True,
                "output_path": output_path,
                "prompt": prompt,
                "duration": duration,
            }
        return {"error": result.get("error", "Sound generation failed")}
    except Exception as e:
        return {"error": str(e)}


# ═══════════════════════════════════════════════════════════════════════════
# Asset Extraction (find_bbox + crop — mirrors Kimi's vision tools)
# ═══════════════════════════════════════════════════════════════════════════

async def _find_bbox_impl(image_path: str, description: str) -> Dict:
    """
    Find bounding box of an element in an image.
    Uses OpenAI Vision API to locate elements.
    """
    try:
        from openai import OpenAI
        import base64

        client = OpenAI()

        # Read and encode image
        with open(image_path, "rb") as f:
            img_data = base64.b64encode(f.read()).decode("utf-8")

        # Determine MIME type
        ext = os.path.splitext(image_path)[1].lower()
        mime = {"png": "image/png", "jpg": "image/jpeg", "jpeg": "image/jpeg", "gif": "image/gif", "webp": "image/webp"}.get(ext.lstrip("."), "image/png")

        response = client.chat.completions.create(
            model="gpt-4.1-mini",
            messages=[
                {
                    "role": "user",
                    "content": [
                        {
                            "type": "text",
                            "text": f'Find the bounding box of "{description}" in this image. Return ONLY a JSON object with keys: x, y, width, height (as percentages 0-100 of image dimensions). If not found, return {{"error": "not found"}}.',
                        },
                        {
                            "type": "image_url",
                            "image_url": {"url": f"data:{mime};base64,{img_data}"},
                        },
                    ],
                }
            ],
            max_tokens=200,
        )

        result_text = response.choices[0].message.content.strip()
        # Extract JSON from response
        json_match = re.search(r'\{[^}]+\}', result_text)
        if json_match:
            bbox = json.loads(json_match.group())
            return {"success": True, "bbox": bbox, "description": description}
        return {"error": "Could not parse bounding box from response"}
    except Exception as e:
        return {"error": str(e)}


async def _crop_image_impl(image_path: str, x: int, y: int, width: int, height: int, output_path: str) -> Dict:
    """Crop an image to the specified region."""
    try:
        from PIL import Image
        img = Image.open(image_path)
        # If percentages, convert to pixels
        img_w, img_h = img.size
        if x <= 100 and y <= 100 and width <= 100 and height <= 100:
            x = int(x / 100 * img_w)
            y = int(y / 100 * img_h)
            width = int(width / 100 * img_w)
            height = int(height / 100 * img_h)

        cropped = img.crop((x, y, x + width, y + height))
        os.makedirs(os.path.dirname(output_path) if os.path.dirname(output_path) else ".", exist_ok=True)
        cropped.save(output_path)
        return {
            "success": True,
            "output_path": output_path,
            "original_size": f"{img_w}x{img_h}",
            "crop_region": f"{x},{y},{width},{height}",
            "cropped_size": f"{cropped.width}x{cropped.height}",
        }
    except Exception as e:
        return {"error": str(e)}


# ═══════════════════════════════════════════════════════════════════════════
# Slides Generator (HTML → PPTX)
# ═══════════════════════════════════════════════════════════════════════════

async def _create_slides_impl(slides_data: List[Dict], output_path: str, theme: str = "default") -> Dict:
    """
    Generate a PPTX presentation from structured slide data.
    Each slide: {"title": "...", "content": "...", "notes": "...", "layout": "..."}
    """
    try:
        code = f"""
import json
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

slides_data = {json.dumps(slides_data)}
output_path = "{output_path}"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Theme colors
themes = {{
    "default": {{"bg": "FFFFFF", "title": "1A1A2E", "text": "333333", "accent": "0066CC"}},
    "dark": {{"bg": "1A1A2E", "title": "FFFFFF", "text": "E0E0E0", "accent": "00D4FF"}},
    "corporate": {{"bg": "F5F5F5", "title": "003366", "text": "444444", "accent": "CC6600"}},
}}
theme = themes.get("{theme}", themes["default"])

for i, slide_data in enumerate(slides_data):
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor.from_string(theme["bg"])

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.7), Inches(1.2))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = slide_data.get("title", f"Slide {{i+1}}")
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor.from_string(theme["title"])

    # Content
    content = slide_data.get("content", "")
    if content:
        content_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(11.7), Inches(4.5))
        tf = content_box.text_frame
        tf.word_wrap = True
        for line in content.split("\\n"):
            p = tf.add_paragraph()
            p.text = line
            p.font.size = Pt(18)
            p.font.color.rgb = RGBColor.from_string(theme["text"])
            p.space_after = Pt(8)

    # Notes
    notes = slide_data.get("notes", "")
    if notes:
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = notes

os.makedirs(os.path.dirname(output_path) if os.path.dirname(output_path) else ".", exist_ok=True)
prs.save(output_path)
print(f"Presentation saved: {{len(slides_data)}} slides, {{os.path.getsize(output_path)}} bytes")
"""
        # First ensure python-pptx is installed
        await _sandbox_mgr.run_code("import subprocess; subprocess.run(['pip', 'install', 'python-pptx', '-q'], capture_output=True)", timeout=30)
        result = await _sandbox_mgr.run_code(code, timeout=30)
        if result.get("success"):
            return {
                "success": True,
                "output_path": output_path,
                "slide_count": len(slides_data),
                "theme": theme,
            }
        return {"error": result.get("error", "Slides generation failed")}
    except Exception as e:
        return {"error": str(e)}


# ═══════════════════════════════════════════════════════════════════════════
# Website Deployment (with HTTP server)
# ═══════════════════════════════════════════════════════════════════════════

async def _deploy_website_impl(local_dir: str, description: str = "") -> Dict:
    """Deploy a static website."""
    if not os.path.exists(os.path.join(local_dir, "index.html")):
        return {"error": "index.html not found in the specified directory"}

    deploy_id = datetime.now().strftime("%Y%m%d_%H%M%S")
    deploy_dir = DEPLOY_DIR / deploy_id
    try:
        shutil.copytree(local_dir, str(deploy_dir))
        return {
            "success": True,
            "deploy_id": deploy_id,
            "deployed_to": str(deploy_dir),
            "description": description,
            "files": os.listdir(str(deploy_dir)),
            "serve_command": f"python -m http.server 8080 --directory {deploy_dir}",
        }
    except Exception as e:
        return {"error": str(e)}


# ═══════════════════════════════════════════════════════════════════════════
# Todo Management (Persistent — .todo.jsonl)
# ═══════════════════════════════════════════════════════════════════════════

def _load_todos() -> List[Dict]:
    """Load todos from persistent storage."""
    if not TODO_FILE.exists():
        return []
    try:
        todos = []
        with open(TODO_FILE, "r") as f:
            for line in f:
                line = line.strip()
                if line:
                    todos.append(json.loads(line))
        return todos
    except Exception:
        return []


def _save_todos(todos: List[Dict]):
    """Save todos to persistent storage."""
    try:
        with open(TODO_FILE, "w") as f:
            for todo in todos:
                f.write(json.dumps(todo) + "\n")
    except Exception as e:
        logger.error(f"Failed to save todos: {e}")


# ═══════════════════════════════════════════════════════════════════════════
# MASTER TOOL ROUTER — All 29 tools
# ═══════════════════════════════════════════════════════════════════════════

async def execute_tool(tool_name: str, arguments: Dict[str, Any]) -> str:
    """Route tool calls to their implementations. All 29 tools."""
    executors = {
        # Task Management
        "todo_read": _todo_read,
        "todo_write": _todo_write,
        # Code Execution
        "ipython": _ipython,
        # File Operations
        "read_file": _read_file,
        "write_file": _write_file,
        "edit_file": _edit_file,
        # Shell
        "shell": _shell,
        # Web Search
        "web_search": _web_search,
        # Browser Automation (7 tools)
        "browser_visit": _browser_visit,
        "browser_click": _browser_click,
        "browser_input": _browser_input,
        "browser_scroll_down": _browser_scroll_down,
        "browser_scroll_up": _browser_scroll_up,
        "browser_screenshot": _browser_screenshot,
        "browser_find": _browser_find,
        # Image Generation
        "generate_image": _generate_image,
        # Speech / Audio
        "generate_speech": _generate_speech,
        "generate_sound_effect": _generate_sound_effect,
        # Data Sources
        "get_data_source": _get_data_source,
        # Asset Extraction
        "find_bbox": _find_bbox,
        "crop_image": _crop_image,
        # Slides
        "create_slides": _create_slides,
        # Website Deployment
        "deploy_website": _deploy_website,
        # Utility
        "get_current_time": _get_current_time,
        "list_workspace": _list_workspace,
        "download_file": _download_file,
        "upload_file": _upload_file,
    }

    executor = executors.get(tool_name)
    if not executor:
        return json.dumps({"error": f"Unknown tool: {tool_name}. Available: {list(executors.keys())}"})

    try:
        result = await executor(arguments)
        return result if isinstance(result, str) else json.dumps(result, default=str)
    except Exception as e:
        logger.error(f"Tool {tool_name} error: {e}", exc_info=True)
        return json.dumps({"error": f"{e.__class__.__name__}: {str(e)}"})


# ═══════════════════════════════════════════════════════════════════════════
# Individual Tool Implementations
# ═══════════════════════════════════════════════════════════════════════════

# ─── 1. Todo Read ─────────────────────────────────────────────────────────

async def _todo_read(args: dict) -> str:
    todos = _load_todos()
    if not todos:
        return json.dumps({"todos": [], "message": "No todos yet."})
    return json.dumps({"todos": todos, "count": len(todos)})


# ─── 2. Todo Write ────────────────────────────────────────────────────────

async def _todo_write(args: dict) -> str:
    todos = args.get("todos", [])
    # Validate: only one in_progress
    in_progress = [t for t in todos if t.get("status") == "in_progress"]
    if len(in_progress) > 1:
        return json.dumps({"error": "Only one task can be in_progress at a time."})
    _save_todos(todos)
    return json.dumps({"success": True, "count": len(todos)})


# ─── 3. IPython ───────────────────────────────────────────────────────────

async def _ipython(args: dict) -> str:
    code = args.get("code", "")
    restart = args.get("restart", False)

    if restart:
        if _sandbox_mgr.use_e2b and _sandbox_mgr.sandbox:
            try:
                await _sandbox_mgr.sandbox.restart_code_context()
            except Exception:
                pass
        else:
            kernel = _get_kernel()
            kernel.restart()
        return json.dumps({"success": True, "output": "Kernel restarted."})

    result = await _sandbox_mgr.run_code(code, timeout=30)
    return json.dumps(result, default=str)


# ─── 4. Read File ─────────────────────────────────────────────────────────

async def _read_file(args: dict) -> str:
    file_path = args.get("file_path", "")
    limit = args.get("limit", 1000)
    offset = args.get("offset", 1)

    if not os.path.isabs(file_path):
        return json.dumps({"error": "file_path must be absolute"})
    if not os.path.exists(file_path):
        return json.dumps({"error": f"File not found: {file_path}"})

    _files_read.add(file_path)

    # Check if binary file (image, video, etc.)
    ext = os.path.splitext(file_path)[1].lower()
    binary_exts = {'.png', '.jpg', '.jpeg', '.gif', '.bmp', '.webp', '.mp4', '.mp3', '.wav', '.pdf', '.zip', '.tar', '.gz'}
    if ext in binary_exts:
        file_size = os.path.getsize(file_path)
        if ext in {'.png', '.jpg', '.jpeg', '.gif', '.bmp', '.webp'}:
            try:
                from PIL import Image
                img = Image.open(file_path)
                return json.dumps({
                    "success": True,
                    "type": "image",
                    "format": img.format,
                    "size": f"{img.width}x{img.height}",
                    "mode": img.mode,
                    "file_size": file_size,
                    "file_path": file_path,
                })
            except Exception:
                pass
        return json.dumps({
            "success": True,
            "type": "binary",
            "extension": ext,
            "file_size": file_size,
            "file_path": file_path,
        })

    try:
        with open(file_path, "r", encoding="utf-8", errors="replace") as f:
            lines = f.readlines()

        start = max(0, offset - 1)
        end = start + limit
        selected = lines[start:end]

        numbered = []
        for i, line in enumerate(selected, start=start + 1):
            truncated = line[:2000] if len(line) > 2000 else line
            numbered.append(f"  {i}\t{truncated.rstrip()}")

        output = "\n".join(numbered)
        if len(output) > 10000:
            output = output[:10000] + "\n... [output truncated]"

        return json.dumps({
            "success": True,
            "content": output,
            "total_lines": len(lines),
            "lines_shown": len(selected),
            "has_more": end < len(lines),
        })
    except Exception as e:
        return json.dumps({"error": str(e)})


# ─── 5. Write File ────────────────────────────────────────────────────────

async def _write_file(args: dict) -> str:
    file_path = args.get("file_path", "")
    content = args.get("content", "")
    append = args.get("append", False)

    if not os.path.isabs(file_path):
        return json.dumps({"error": "file_path must be absolute"})
    if len(content) > 100000:
        return json.dumps({"error": "Content exceeds 100000 character limit"})

    try:
        os.makedirs(os.path.dirname(file_path), exist_ok=True)
        mode = "a" if append else "w"
        with open(file_path, mode, encoding="utf-8") as f:
            f.write(content)
        return json.dumps({
            "success": True,
            "file_path": file_path,
            "bytes_written": len(content.encode("utf-8")),
            "mode": "append" if append else "overwrite",
        })
    except Exception as e:
        return json.dumps({"error": str(e)})


# ─── 6. Edit File ─────────────────────────────────────────────────────────

async def _edit_file(args: dict) -> str:
    file_path = args.get("file_path", "")
    old_string = args.get("old_string", "")
    new_string = args.get("new_string", "")
    replace_all = args.get("replace_all", False)

    if file_path not in _files_read:
        return json.dumps({"error": "Must read_file before editing. Read the file first."})
    if old_string == new_string:
        return json.dumps({"error": "old_string and new_string must be different"})

    try:
        with open(file_path, "r", encoding="utf-8") as f:
            content = f.read()

        count = content.count(old_string)
        if count == 0:
            return json.dumps({"error": "old_string not found in file"})
        if count > 1 and not replace_all:
            return json.dumps({
                "error": f"old_string found {count} times. Use replace_all: true or provide more unique context."
            })

        new_content = content.replace(old_string, new_string) if replace_all else content.replace(old_string, new_string, 1)

        with open(file_path, "w", encoding="utf-8") as f:
            f.write(new_content)

        return json.dumps({
            "success": True,
            "replacements": count if replace_all else 1,
        })
    except Exception as e:
        return json.dumps({"error": str(e)})


# ─── 7. Shell ─────────────────────────────────────────────────────────────

async def _shell(args: dict) -> str:
    command = args.get("command", "")
    timeout = args.get("timeout", 600)
    result = await _sandbox_mgr.run_shell(command, timeout)
    return json.dumps(result, default=str)


# ─── 8. Web Search ────────────────────────────────────────────────────────

async def _web_search(args: dict) -> str:
    query = args.get("query", "")
    count = args.get("count", 5)
    result = await _perform_web_search(query, count)
    return json.dumps(result, default=str)


# ─── 9-15. Browser Tools ─────────────────────────────────────────────────

async def _browser_visit(args: dict) -> str:
    url = args.get("url", "")
    need_screenshot = args.get("need_screenshot", False)
    result = await _browser.visit(url, need_screenshot)
    return json.dumps(result, default=str)


async def _browser_click(args: dict) -> str:
    element_index = args.get("element_index", 0)
    result = await _browser.click(element_index)
    return json.dumps(result, default=str)


async def _browser_input(args: dict) -> str:
    element_index = args.get("element_index", 0)
    content = args.get("content", "")
    result = await _browser.input_text(element_index, content)
    return json.dumps(result, default=str)


async def _browser_scroll_down(args: dict) -> str:
    amount = args.get("scroll_amount", 500)
    result = await _browser.scroll("down", amount)
    return json.dumps(result, default=str)


async def _browser_scroll_up(args: dict) -> str:
    amount = args.get("scroll_amount", 500)
    result = await _browser.scroll("up", amount)
    return json.dumps(result, default=str)


async def _browser_screenshot(args: dict) -> str:
    download_path = args.get("download_path")
    result = await _browser.screenshot(download_path)
    return json.dumps(result, default=str)


async def _browser_find(args: dict) -> str:
    keyword = args.get("keyword", "")
    skip = args.get("skip", 0)
    result = await _browser.find_text(keyword, skip)
    return json.dumps(result, default=str)


# ─── 16. Generate Image ──────────────────────────────────────────────────

async def _generate_image(args: dict) -> str:
    prompt = args.get("prompt", "")
    output_path = args.get("output_path", str(OUTPUT_DIR / f"image_{int(time.time())}.png"))
    ratio = args.get("ratio", "1:1")
    resolution = args.get("resolution", "1K")
    result = await _generate_image_impl(prompt, output_path, ratio, resolution)
    return json.dumps(result, default=str)


# ─── 17. Generate Speech ─────────────────────────────────────────────────

async def _generate_speech(args: dict) -> str:
    text = args.get("text", "")
    output_path = args.get("output_path", str(OUTPUT_DIR / f"speech_{int(time.time())}.mp3"))
    voice = args.get("voice", "alloy")
    speed = args.get("speed", 1.0)
    result = await _generate_speech_impl(text, output_path, voice, speed)
    return json.dumps(result, default=str)


# ─── 18. Generate Sound Effect ────────────────────────────────────────────

async def _generate_sound_effect(args: dict) -> str:
    prompt = args.get("prompt", "")
    output_path = args.get("output_path", str(OUTPUT_DIR / f"sfx_{int(time.time())}.wav"))
    duration = args.get("duration", 3.0)
    result = await _generate_sound_effect_impl(prompt, output_path, duration)
    return json.dumps(result, default=str)


# ─── 19. Get Data Source ──────────────────────────────────────────────────

async def _get_data_source(args: dict) -> str:
    source = args.get("source", "")
    query = args.get("query", "")
    params = args.get("params", {})

    source_map = {
        "yahoo_finance": _data_client.yahoo_finance,
        "world_bank": _data_client.world_bank,
        "world_bank_open_data": _data_client.world_bank,
        "arxiv": _data_client.arxiv_search,
        "google_scholar": _data_client.google_scholar,
    }

    handler = source_map.get(source)
    if not handler:
        return json.dumps({
            "error": f"Unknown data source: {source}",
            "available": list(source_map.keys()),
        })

    if source in ("world_bank", "world_bank_open_data"):
        result = await handler(query, params.get("country", "all"), params)
    else:
        result = await handler(query, **({"max_results": params.get("max_results", 10)} if source in ("arxiv", "google_scholar") else {}))

    return json.dumps(result, default=str)


# ─── 20-21. Asset Extraction ─────────────────────────────────────────────

async def _find_bbox(args: dict) -> str:
    image_path = args.get("image_path", "")
    description = args.get("description", "")
    result = await _find_bbox_impl(image_path, description)
    return json.dumps(result, default=str)


async def _crop_image(args: dict) -> str:
    image_path = args.get("image_path", "")
    x = args.get("x", 0)
    y = args.get("y", 0)
    width = args.get("width", 100)
    height = args.get("height", 100)
    output_path = args.get("output_path", str(OUTPUT_DIR / f"cropped_{int(time.time())}.png"))
    result = await _crop_image_impl(image_path, x, y, width, height, output_path)
    return json.dumps(result, default=str)


# ─── 22. Create Slides ───────────────────────────────────────────────────

async def _create_slides(args: dict) -> str:
    slides_data = args.get("slides", [])
    output_path = args.get("output_path", str(OUTPUT_DIR / f"presentation_{int(time.time())}.pptx"))
    theme = args.get("theme", "default")
    result = await _create_slides_impl(slides_data, output_path, theme)
    return json.dumps(result, default=str)


# ─── 23. Deploy Website ──────────────────────────────────────────────────

async def _deploy_website(args: dict) -> str:
    local_dir = args.get("local_dir", "")
    description = args.get("description", "")
    result = await _deploy_website_impl(local_dir, description)
    return json.dumps(result, default=str)


# ─── 24. Get Current Time ────────────────────────────────────────────────

async def _get_current_time(args: dict) -> str:
    now = datetime.now()
    return json.dumps({
        "datetime": now.isoformat(),
        "date": now.strftime("%Y-%m-%d"),
        "time": now.strftime("%H:%M:%S"),
        "timestamp": int(now.timestamp()),
        "timezone": "UTC",
    })


# ─── 25. List Workspace ──────────────────────────────────────────────────

async def _list_workspace(args: dict) -> str:
    path = args.get("path", str(WORKSPACE))
    try:
        entries = []
        for entry in sorted(os.listdir(path)):
            full_path = os.path.join(path, entry)
            stat = os.stat(full_path)
            entries.append({
                "name": entry,
                "type": "directory" if os.path.isdir(full_path) else "file",
                "size": stat.st_size,
                "modified": datetime.fromtimestamp(stat.st_mtime).isoformat(),
            })
        return json.dumps({"success": True, "path": path, "entries": entries})
    except Exception as e:
        return json.dumps({"error": str(e)})


# ─── 26. Download File ───────────────────────────────────────────────────

async def _download_file(args: dict) -> str:
    url = args.get("url", "")
    output_path = args.get("output_path", "")
    if not output_path:
        filename = url.split("/")[-1].split("?")[0] or "download"
        output_path = str(OUTPUT_DIR / filename)
    try:
        import requests
        response = requests.get(url, timeout=60, stream=True)
        response.raise_for_status()
        os.makedirs(os.path.dirname(output_path), exist_ok=True)
        with open(output_path, "wb") as f:
            for chunk in response.iter_content(chunk_size=8192):
                f.write(chunk)
        return json.dumps({
            "success": True,
            "output_path": output_path,
            "size": os.path.getsize(output_path),
            "content_type": response.headers.get("content-type", "unknown"),
        })
    except Exception as e:
        return json.dumps({"error": str(e)})


# ─── 27. Upload File ─────────────────────────────────────────────────────

async def _upload_file(args: dict) -> str:
    file_path = args.get("file_path", "")
    if not os.path.exists(file_path):
        return json.dumps({"error": f"File not found: {file_path}"})
    try:
        # Copy to upload directory with unique name
        filename = os.path.basename(file_path)
        dest = str(UPLOAD_DIR / f"{int(time.time())}_{filename}")
        shutil.copy2(file_path, dest)
        return json.dumps({
            "success": True,
            "original_path": file_path,
            "uploaded_path": dest,
            "size": os.path.getsize(dest),
        })
    except Exception as e:
        return json.dumps({"error": str(e)})

# ═══════════════════════════════════════════════════════════════════════════
# ToolExecutor Class (Orchestrator Interface)
# ═══════════════════════════════════════════════════════════════════════════

class ToolExecutor:
    """
    Persistent executor for a single agent session.
    Maintains state for browser, sandbox, and other tools.
    """

    def __init__(self):
        self.browser = PlaywrightBrowser()
        self.sandbox_mgr = _sandbox_mgr
        self.files_read = set()

    async def execute(self, tool_name: str, arguments: Dict[str, Any]) -> str:
        """Execute a tool and return the result as a string."""
        # Update global state for this call (since implementations use globals)
        global _browser, _files_read
        _browser = self.browser
        _files_read = self.files_read
        
        return await execute_tool(tool_name, arguments)

    async def cleanup(self):
        """Cleanup persistent resources."""
        if self.browser:
            try:
                if self.browser.page:
                    await self.browser.page.close()
                if self.browser.browser:
                    await self.browser.browser.close()
                if self.browser.playwright:
                    await self.browser.playwright.stop()
            except Exception:
                pass
        # Note: Sandbox is global in this implementation, but could be per-session
